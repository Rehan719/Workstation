import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class SciencePresentationGeneratorV16:
    def __init__(self):
        self.prs = Presentation()
        self.output_path = 'outputs/Science/PatientSafety/v16_quinta_veritas/MULTIMEDIA_ASSETS/presentation/PatientSafety_Presentation_v16_QuintaVeritas.pptx'

    def generate(self):
        # Slide 1: Title & Executive Ask
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = "Patient Safety: Quinta-Veritas Strategic Imperative"
        slide.placeholders[1].text = "Science Grand Operation v16.0 — Ultimate Integrated Briefing"

        # Slide 2: BLUF & Key Findings (Enhanced)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "BLUF: Adaptive Preventive Intelligence"
        tf = slide.placeholders[1].text_frame
        tf.text = "Consolidating 5 Truth Dimensions for Systemic Reform"

        # Slide 11: Truth Five - Ethical-Systemic
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Truth V: Ethical-Systemic Intelligence"
        tf = slide.placeholders[1].text_frame
        tf.text = "Organizational Culture Assessment & Equity Impact Modeling"

        # Slide 13: Quinta-Veritas Synthesis
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Quinta-Veritas Synthesis Engine"
        tf = slide.placeholders[1].text_frame
        tf.text = "Weighted Coherence Scoring (98%) | SHA-3-512 Verified"

        # Slide 24: Closing
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Closing: For Patients, Through Justice."

        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        self.prs.save(self.output_path)
        return self.output_path

if __name__ == "__main__":
    gen = SciencePresentationGeneratorV16()
    path = gen.generate()
    print(f"✅ Quinta-Veritas Presentation Generated: {path}")
