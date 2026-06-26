import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class SciencePresentationGeneratorV14:
    def __init__(self, audience='ELT'):
        self.prs = Presentation()
        self.audience = audience
        self.output_path = 'outputs/v14/science/presentation/PatientSafety_Presentation_v14.pptx'

    def _add_video_placeholder(self, slide, video_filename, position):
        x, y, width, height = position
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.text_frame.text = f"EMBED VIDEO: {video_filename}\nAuto-play on slide advance"

    def _add_nested_infographic_placeholder(self, slide, infographic_id, position):
        x, y, width, height = position
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(220, 230, 240)
        placeholder.text_frame.text = f"NESTED INFOGRAPHIC: {infographic_id}\nInteractive Visualisation"

    def calculate_temporal_convergence_v14(self):
        """Implement context-aware temporal weighting logic"""
        weights = {'Truth_I': 0.30, 'Truth_II': 0.25, 'Truth_III': 0.25, 'Truth_IV': 0.20}
        if self.audience == 'ELT':
            weights['Truth_IV'] += 0.10
            weights['Truth_II'] -= 0.05
        elif self.audience == 'Board_QC':
            weights['Truth_III'] += 0.15
            weights['Truth_I'] -= 0.10
        return weights

    def generate(self):
        # 1-4: Executive Hook
        self._add_title_slide()
        self._add_bluf_slide()
        self._add_scientific_validation_slide()
        self._add_opportunity_window_slide()

        # 5-10: Mushahida
        self._add_trial_design_flaws_slide()
        self._add_regulatory_status_quo_slide()
        self._add_evidence_base_slide()
        self._add_des_tragedy_slide()
        self._add_stakeholder_mapping_slide()
        self._add_cdmo_risk_slide()

        # 11-16: Jaiza
        self._add_proceduralism_trap_slide()
        self._add_trend_recognition_slide()
        self._add_five_point_framework_slide()
        self._add_strategic_options_slide()
        self._add_risk_benefit_assessment_slide()
        self._add_convergence_analysis_slide()

        # 17-24: Muaina
        self._add_strategic_position_slide()
        self._add_regulatory_affairs_proposal_slide()
        self._add_commercial_proposal_slide()
        self._add_rd_proposal_slide()
        self._add_implementation_roadmap_slide()
        self._add_budget_roi_slide()
        self._add_escalation_request_slide()
        self._add_closing_slide()

        self.prs.save(self.output_path)
        return self.output_path

    # --- SLIDE DEFINITIONS ---
    def _add_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = "Patient Safety Intelligence: Quadra-Veritas Strategic Imperative"
        slide.placeholders[1].text = "Science Grand Operation v14.0 — Multimedia Intelligence Briefing"

    def _add_bluf_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "The Patient Safety Imperative (BLUF)"
        tf = slide.placeholders[1].text_frame
        tf.text = "Two Critical Gaps. One Strategic Opportunity."
        self._add_video_placeholder(slide, "patient_impact_vignette.mp4", (Inches(7), Inches(5), Inches(2.5), Inches(1.5)))

    def _add_scientific_validation_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Definitive Scientific Validation (2025-2026)"
        self._add_video_placeholder(slide, "aav_germ_cell_explainer.mp4", (Inches(7), Inches(5), Inches(2.5), Inches(1.5)))

    def _add_opportunity_window_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "The Opportunity Window"

    def _add_trial_design_flaws_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Therapeutic Landscape & Trial Design Flaws"

    def _add_regulatory_status_quo_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Regulatory Status Quo Analysis"

    def _add_evidence_base_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Comprehensive Scientific Evidence Base"

    def _add_des_tragedy_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Historical Precedent: The DES Tragedy"

    def _add_stakeholder_mapping_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Stakeholder Mapping & Patient Impact"

    def _add_cdmo_risk_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Why This Matters to Lonza as a CDMO"

    def _add_proceduralism_trap_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "The Proceduralism Trap"

    def _add_trend_recognition_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Pattern Recognition: Trends & Disruptors"

    def _add_five_point_framework_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Recommended Actions: The Five-Point Framework"

    def _add_strategic_options_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Strategic Options Analysis"

    def _add_risk_benefit_assessment_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Risk-Benefit Assessment (Patient-Centred)"

    def _add_convergence_analysis_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Quadra-Veritas Convergence Analysis"

    def _add_strategic_position_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Alignment with Lonza's Strategic Position"

    def _add_regulatory_affairs_proposal_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Proposal for Regulatory Affairs"

    def _add_commercial_proposal_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Proposal for Commercial Development"

    def _add_rd_proposal_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Proposal for R&D"

    def _add_implementation_roadmap_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Integrated Implementation Roadmap"

    def _add_budget_roi_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Consolidated Budget & ROI Analysis"

    def _add_escalation_request_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Request for Escalation & Next Steps"

    def _add_closing_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Closing: Upholding Our Commitment to Patients"

if __name__ == "__main__":
    gen = SciencePresentationGeneratorV14(audience='ELT')
    pptx_path = gen.generate()
    print(f"✅ Full 24-slide presentation generated: {pptx_path}")
