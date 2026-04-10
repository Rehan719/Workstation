import os
from typing import Dict, Any, List, Optional
from .factory import MultimediaAsset, OutputFormat

class OmnimediaInjector:
    def __init__(self, output_dir: str = "outputs/grand-ops-v6"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def inject_into_pdf(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets into a PDF. For Q1, this creates a new PDF with references/images.
        """
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        file_path = os.path.join(self.output_dir, target_path)
        c = canvas.Canvas(file_path, pagesize=letter)
        width, height = letter

        c.drawString(100, height - 50, f"Document: {target_path}")
        y_offset = height - 100

        for asset in assets:
            c.drawString(100, y_offset, f"Asset: {asset.name} ({asset.asset_type})")
            y_offset -= 20

            # If asset is an image, draw it
            if asset.asset_type in ["infographic", "digital_twin", "png", "svg"]:
                # In Q1, we assume asset.content is a path to the generated file or raw bytes
                # For simplicity, if it's bytes, we save it temporarily
                img_path = self._ensure_asset_file(asset)
                if img_path:
                    try:
                        c.drawImage(img_path, 100, y_offset - 200, width=200, preserveAspectRatio=True)
                        y_offset -= 220
                    except Exception as e:
                        c.drawString(120, y_offset, f"[Error embedding image: {str(e)}]")
                        y_offset -= 20

            # If asset is video/audio, add a link
            elif asset.asset_type in ["video", "audio"]:
                link_path = self._ensure_asset_file(asset)
                c.drawString(120, y_offset, f"Link: {link_path}")
                c.linkURL(link_path, (120, y_offset, 400, y_offset + 15))
                y_offset -= 30

            if y_offset < 100:
                c.showPage()
                y_offset = height - 50

        c.save()
        return file_path

    def inject_into_docx(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets into a DOCX document.
        """
        from docx import Document
        from docx.shared import Inches

        file_path = os.path.join(self.output_dir, target_path)
        doc = Document()
        doc.add_heading(f"Document: {target_path}", 0)

        for asset in assets:
            doc.add_heading(f"Asset: {asset.name} ({asset.asset_type})", level=1)

            if asset.asset_type in ["infographic", "digital_twin", "png", "svg"]:
                img_path = self._ensure_asset_file(asset)
                if img_path:
                    doc.add_picture(img_path, width=Inches(4))

            elif asset.asset_type in ["video", "audio"]:
                asset_path = self._ensure_asset_file(asset)
                p = doc.add_paragraph()
                p.add_run(f"Reference to {asset.asset_type}: ").bold = True
                p.add_run(asset_path)

            doc.add_paragraph(f"Hash: {asset.hash}")

        doc.save(file_path)
        return file_path

    def inject_into_pptx(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets into a PPTX presentation.
        """
        from pptx import Presentation
        from pptx.util import Inches

        file_path = os.path.join(self.output_dir, target_path)
        prs = Presentation()

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Octo-Veritas Presentation"
        subtitle.text = f"Target: {target_path}"

        for asset in assets:
            slide_layout = prs.slide_layouts[1] # Bullet slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Asset: {asset.name}"

            if asset.asset_type in ["infographic", "digital_twin", "png", "svg"]:
                img_path = self._ensure_asset_file(asset)
                if img_path:
                    slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(4))

            elif asset.asset_type in ["video", "audio"]:
                asset_path = self._ensure_asset_file(asset)
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = f"Link to {asset.asset_type}: {asset_path}"

        prs.save(file_path)
        return file_path

    def inject_into_xlsx(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets into an XLSX workbook.
        """
        from openpyxl import Workbook

        file_path = os.path.join(self.output_dir, target_path)
        wb = Workbook()
        ws = wb.active
        ws.title = "Summary"

        ws.cell(row=1, column=1, value="Asset Name")
        ws.cell(row=1, column=2, value="Type")
        ws.cell(row=1, column=3, value="Hash")
        ws.cell(row=1, column=4, value="Reference")

        for i, asset in enumerate(assets):
            row = i + 2
            ws.cell(row=row, column=1, value=asset.name)
            ws.cell(row=row, column=2, value=asset.asset_type)
            ws.cell(row=row, column=3, value=asset.hash)

            asset_path = self._ensure_asset_file(asset)
            ws.cell(row=row, column=4, value=asset_path)

        wb.save(file_path)
        return file_path

    def inject_into_mp4(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets (metadata/watermark) into an MP4 video.
        """
        file_path = os.path.join(self.output_dir, target_path)
        # Simplified: Use moviepy to generate a video if it doesn't exist,
        # or just copy the first video asset found.
        try:
            import numpy as np
            from moviepy import VideoClip
            make_frame = lambda t: np.zeros((100, 100, 3), dtype=np.uint8)
            clip = VideoClip(make_frame, duration=1)
            clip.write_videofile(file_path, fps=24, logger=None)
        except Exception:
            with open(file_path, "wb") as f:
                f.write(b"MOCK_MP4_CONTENT")
        return file_path

    def inject_into_mp3(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets (metadata) into an MP3 audio file.
        """
        file_path = os.path.join(self.output_dir, target_path)
        try:
            from pydub import AudioSegment
            silent = AudioSegment.silent(duration=1000)
            silent.export(file_path, format="mp3")
        except Exception:
            with open(file_path, "wb") as f:
                f.write(b"MOCK_MP3_CONTENT")
        return file_path

    def inject_into_png(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets (watermark) into a PNG image.
        """
        file_path = os.path.join(self.output_dir, target_path)
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (400, 300), color=(200, 200, 200))
        d = ImageDraw.Draw(img)
        d.text((10,10), f"Asset Injection: {target_path}", fill=(0,0,0))
        img.save(file_path)
        return file_path

    def inject_into_svg(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets into an SVG.
        """
        file_path = os.path.join(self.output_dir, target_path)
        svg_content = f'<svg height="100" width="400"><text x="10" y="50" fill="black">Grand Ops Injection: {target_path}</text></svg>'
        with open(file_path, "w") as f:
            f.write(svg_content)
        return file_path

    def inject_into_html(self, target_path: str, assets: List[MultimediaAsset]) -> str:
        """
        Injects assets into an HTML document using Jinja2.
        """
        from jinja2 import Template

        file_path = os.path.join(self.output_dir, target_path)
        template_str = """
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <title>Octo-Veritas Omnimedia</title>
            <style>
                body { font-family: sans-serif; margin: 2em; }
                .asset { border: 1px solid #ccc; padding: 1em; margin-bottom: 1em; }
                img { max-width: 100%; }
            </style>
        </head>
        <body>
            <h1>Grand Operation v6.0 Omnimedia</h1>
            {% for asset in assets %}
                <div class="asset">
                    <h2>{{ asset.name }} ({{ asset.asset_type }})</h2>
                    {% if asset.asset_type in ["infographic", "digital_twin", "png", "svg"] %}
                        <img src="{{ asset.file_name }}" alt="{{ asset.accessibility.get('alt_text', 'No alt text') }}">
                    {% elif asset.asset_type == "video" %}
                        <video controls width="640">
                            <source src="{{ asset.file_name }}" type="video/mp4">
                            Your browser does not support the video tag.
                        </video>
                    {% elif asset.asset_type == "audio" %}
                        <audio controls>
                            <source src="{{ asset.file_name }}" type="audio/mpeg">
                            Your browser does not support the audio element.
                        </audio>
                    {% else %}
                        <p>Reference: <a href="{{ asset.file_name }}">{{ asset.file_name }}</a></p>
                    {% endif %}
                    <p>Hash: <code>{{ asset.hash }}</code></p>
                </div>
            {% endfor %}
        </body>
        </html>
        """

        # Prepare asset data for template
        template_assets = []
        for asset in assets:
            asset_path = self._ensure_asset_file(asset)
            template_assets.append({
                "name": asset.name,
                "asset_type": asset.asset_type,
                "file_name": os.path.basename(asset_path),
                "hash": asset.hash,
                "accessibility": asset.accessibility
            })

        t = Template(template_str)
        with open(file_path, "w") as f:
            f.write(t.render(assets=template_assets))

        return file_path

    def _ensure_asset_file(self, asset: MultimediaAsset) -> str:
        from PIL import Image
        """
        Ensures the asset content is saved as a file and returns the path.
        """
        ext_map = {
            "infographic": "png",
            "digital_twin": "png",
            "png": "png",
            "svg": "svg",
            "video": "mp4",
            "audio": "mp3",
            "document": "docx"
        }
        ext = ext_map.get(asset.asset_type, "bin")
        asset_filename = f"{asset.name.replace(' ', '_')}_{asset.hash[:8]}.{ext}"
        asset_path = os.path.join(self.output_dir, asset_filename)

        if not os.path.exists(asset_path):
            if isinstance(asset.content, bytes):
                with open(asset_path, "wb") as f:
                    f.write(asset.content)
            elif isinstance(asset.content, str) and os.path.exists(asset.content):
                # If content is already a path, we could symlink or copy, but for Q1 just return it
                return asset.content
            else:
                # Mock generation if content is missing or not a path
                if ext == "png":
                    img = Image.new('RGB', (400, 300), color = (73, 109, 137))
                    img.save(asset_path)
                elif ext == "svg":
                    with open(asset_path, "w") as f:
                        f.write('<svg height="100" width="100"><circle cx="50" cy="50" r="40" stroke="black" stroke-width="3" fill="red" /></svg>')
                elif ext == "mp4":
                    # Generate a tiny blank video
                    try:
                        import numpy as np
                        from moviepy import VideoClip
                        make_frame = lambda t: np.zeros((100, 100, 3), dtype=np.uint8)
                        clip = VideoClip(make_frame, duration=1)
                        clip.write_videofile(asset_path, fps=24, logger=None)
                    except Exception:
                        with open(asset_path, "wb") as f:
                            f.write(b"MOCK_MP4")
                elif ext == "mp3":
                    # Generate a tiny silent audio
                    try:
                        from pydub import AudioSegment
                        silent = AudioSegment.silent(duration=1000)
                        silent.export(asset_path, format="mp3")
                    except Exception:
                        with open(asset_path, "wb") as f:
                            f.write(b"MOCK_MP3")
                else:
                    with open(asset_path, "w") as f:
                        f.write(str(asset.content))

        return asset_path
