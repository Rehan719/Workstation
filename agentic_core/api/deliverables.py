"""
Living Deliverables — the platform's outputs as first-class, reconfigurable, re-runnable artefacts.

A deliverable (report · presentation · website · app spec · service spec · brief) is produced on
Workstation's OWN native fabric (in-house-first, with served_by provenance), persisted, and kept
LIVING: it can be re-generated and reconfigured (different brief / sections) at any time, keeping a
version history. Deliverables can be grounded in a live VSB entity and are filed under that VSB.

  GET  /api/v1/deliverables/types              — the deliverable catalogue (type → default sections)
  POST /api/v1/deliverables/produce            — produce a new living deliverable on the native fabric
  GET  /api/v1/deliverables                     — list (optionally ?vsb_id=)
  GET  /api/v1/deliverables/{id}                — one deliverable + its version history
  POST /api/v1/deliverables/{id}/regenerate     — re-run / reconfigure (appends a new version)
"""
from __future__ import annotations

import html as _htmlmod
import json
import re
import time
import uuid
from pathlib import Path
from agentic_core.config import atomic_write_json, data_path
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel

from agentic_core.auth.core import get_current_user, request_owner_id, user_can_access


def _require_deliverable_access(d: Dict[str, Any], user: dict | None) -> None:
    """§14×§13 (W320) — a living deliverable is its owner's confidential work product: previously
    the store had NO owner concept, so any user could list, read, export, and refine any tenant's
    deliverables. 404 (never 403) when scoped out; legacy/unowned records are admin-only under
    auth; single-user mode (auth off) stays unguarded."""
    u = user if isinstance(user, dict) else None
    if not user_can_access(u, d.get("owner_id")):
        raise HTTPException(status_code=404, detail=f"Deliverable {d.get('id')} not found.")

from agentic_core.ai.gateway import gateway
from agentic_core.vbs.quality import assure_delivery

router = APIRouter(prefix="/api/v1/deliverables", tags=["living-deliverables"])

_STORE = data_path("deliverables.json")

# Each living-deliverable type maps to a sensible default section structure (reconfigurable).
_TYPES: Dict[str, List[str]] = {
    "report": ["Executive Summary", "Background", "Analysis", "Findings", "Recommendations", "Next Steps"],
    "presentation": ["Title & Hook", "Problem", "Solution", "Market", "Business Model", "Traction", "The Ask"],
    "website": ["Hero", "Value Proposition", "Features", "How It Works", "Pricing", "Call To Action"],
    "app": ["Overview", "Core Features", "User Flows", "Solution Architecture", "MVP Scope", "Roadmap"],
    "service": ["Service Overview", "Offering", "Delivery Model", "Pricing", "SLA & Quality", "Next Steps"],
    "brief": ["Objective", "Context", "Key Factors", "Approach", "Next Steps"],
}


def _load() -> List[Dict[str, Any]]:
    if _STORE.exists():
        try:
            return json.loads(_STORE.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return []
    return []


def _save(rows: List[Dict[str, Any]]) -> None:
    _STORE.parent.mkdir(parents=True, exist_ok=True)
    atomic_write_json(_STORE, rows[-300:])


def _grounding(vsb_id: Optional[str]) -> str:
    if not vsb_id:
        return ""
    try:
        from agentic_core.api.vsb import _load_vsb
        v = _load_vsb(vsb_id)
        if not v:
            return ""
        return (f"\nGround this deliverable in the live VSB: {v.get('name')} "
                f"(domain: {v.get('domain')}; mission: {v.get('challenge', '')}).")
    except Exception:
        return ""


def _weakest_sections(content: str, secs: List[str]) -> List[str]:
    """The declared sections a draft does NOT actually cover — the concrete thing an improvement
    pass should attack. Measured from the draft, never guessed."""
    low = (content or "").lower()
    return [s for s in secs if str(s).lower() not in low]


async def _generate(d_type: str, title: str, brief: str, domain: str,
                    vsb_id: Optional[str], sections: List[str],
                    prior: str = "") -> Dict[str, Any]:
    secs = sections or _TYPES.get(d_type, _TYPES["report"])
    prompt = (
        f"Produce a {d_type} titled '{title or brief[:60]}'.\n"
        f"Brief: {brief}\nDomain: {domain}{_grounding(vsb_id)}\n\n"
        + "\n".join(f"## {s}" for s in secs)
    )
    # §13 (W425) — a regeneration used to be a fresh roll of the dice: the model never saw the draft
    # it was replacing, so "every deliverable is ALIVE" meant only that a version record existed and
    # the user could press the button again. An improvement pass now receives the prior draft AND the
    # sections it MEASURABLY fails to cover, so it has something specific to fix rather than a hope.
    if prior.strip():
        missing = _weakest_sections(prior, secs)
        focus = ("Sections the current draft does NOT cover — address these first: "
                 + ", ".join(missing) + "\n\n") if missing else (
                 "Every declared section is present; deepen the weakest evidence and specificity.\n\n")
        prompt = (
            f"IMPROVE the existing {d_type} below. Keep what is already strong; do not restart.\n"
            f"Brief: {brief}\nDomain: {domain}{_grounding(vsb_id)}\n\n"
            + focus
            + "Return the FULL improved document under these headings:\n"
            + "\n".join(f"## {s}" for s in secs)
            + "\n\n--- CURRENT DRAFT ---\n" + prior[:6000]
        )
    meta = await gateway.query_meta(prompt, agent=f"deliverable:{d_type}", timeout=30.0,
                                    augment=False)   # W332 — deliverables are saved+exported: no cross-request recall
    return {
        "content": meta.get("output", ""),
        "sections": secs,
        # W425 — reported by the function that actually built the prompt, NOT inferred by the caller
        # from whether it MEANT to pass a prior draft. A first cut set this flag from the caller's own
        # variable, so it kept claiming "compared against the prior draft" after the prior stopped
        # being passed at all — a self-attesting flag, which is the §10 defect this codebase just
        # spent W419 removing. Derived from the prompt that was really sent.
        "improved_from_prior": bool(prior.strip()),
        "improvement_focus": _weakest_sections(prior, secs) if prior.strip() else [],
        "ai_provenance": {
            "posture": "in-house-first",
            "served_by": meta.get("served_by", "native"),
            "is_external": bool(meta.get("is_external")),
        },
    }


class ProduceRequest(BaseModel):
    type: str = "report"
    title: str = ""
    brief: str
    domain: str = "enterprise"
    vsb_id: Optional[str] = None
    sections: List[str] = []          # optional override (reconfigure the structure)
    content: str = ""                 # §4.9 (W306): verbatim ingest — already-produced work
                                      # becomes a living deliverable WITHOUT regeneration


class RegenerateRequest(BaseModel):
    brief: Optional[str] = None       # reconfigure the brief …
    sections: Optional[List[str]] = None  # … and/or the section structure
    content: str = ""                 # §3A (W308): DEVELOP — refined text persists as the next
                                      # version VERBATIM (no regeneration), same QMS gate


@router.get("/types")
async def deliverable_types():
    return {"types": [{"id": k, "sections": v} for k, v in _TYPES.items()],
            "posture": "in-house-first"}


@router.get("/output-formats")
async def output_formats():
    """Output formats a deliverable can be exported in (§4.9). Honest by construction: `live` formats are
    rendered end-to-end IN-HOUSE today (real deterministic renders of the deliverable's own content) via
    `/export?format=`; `catalogue` formats are the richer omnimedia targets (pptx/pdf/docx/mp4/…) NOT yet
    produced — listed, never faked."""
    live = [
        {"id": "md", "label": "Markdown", "kind": "document"},
        {"id": "html", "label": "HTML document / website", "kind": "website"},
        {"id": "slides", "label": "HTML presentation (slide deck)", "kind": "presentation"},
        {"id": "txt", "label": "Plain text", "kind": "document"},
        {"id": "json", "label": "Structured JSON", "kind": "data"},
        {"id": "video-html", "label": "Self-playing video (animated HTML — auto-advancing scenes)", "kind": "video"},
    ]
    if _PDF_OK:   # produced in-house by fpdf2 when available
        live.append({"id": "pdf", "label": "PDF document", "kind": "document"})
    if _DOCX_OK:  # produced in-house by python-docx when available
        live.append({"id": "docx", "label": "Word document (.docx)", "kind": "document"})
    if _PPTX_OK:  # produced in-house by python-pptx when available
        live.append({"id": "pptx", "label": "PowerPoint (.pptx)", "kind": "presentation"})
    if _XLSX_OK:  # produced in-house by openpyxl when available
        live.append({"id": "xlsx", "label": "Excel spreadsheet (.xlsx)", "kind": "data"})
    # W372 — svg is dependency-free; png needs Pillow. mp4/mp3 stay in the catalogue: there is no
    # media encoder here, and a "video" that is silently a slideshow would be a fabrication.
    live.append({"id": "svg", "label": "Vector summary card (.svg)", "kind": "image"})
    if _PNG_OK:
        live.append({"id": "png", "label": "Image summary card (.png)", "kind": "image"})
    try:
        from agentic_core.omnimedia.factory import OutputFormat
        catalogue = [f.value for f in OutputFormat if f.value not in {x["id"] for x in live}]
    except Exception:
        catalogue = []
    return {
        "live": live,
        "live_ids": [x["id"] for x in live],
        "catalogue_not_yet_produced": catalogue,
        "source": "agentic_core.omnimedia (the platform's own multimedia output factory)",
        "note": "live = real in-house renders via /export?format=; catalogue = documented targets, not faked.",
    }


@router.post("/produce")
async def produce(req: ProduceRequest, user: dict | None = Depends(get_current_user)):
    """Produce a new LIVING deliverable on the native fabric (in-house provenance)."""
    if req.type not in _TYPES and not req.sections:
        raise HTTPException(status_code=400,
                            detail=f"Unknown type '{req.type}'. Known: {list(_TYPES)} (or pass sections).")
    _t0 = time.time()
    if req.content.strip():
        # §4.9 (W306) — verbatim ingest: output already produced elsewhere on the platform (a Genesis
        # journey, a domain-tool run) becomes a LIVING deliverable — carried verbatim (no regeneration,
        # honest provenance) so every live export format is reachable for it. Still gated below by the
        # same §10/§11 assure_delivery as generated content.
        _secs = req.sections or [m.group(1).strip() for m in
                                 re.finditer(r"^#{1,2}\s+(.+)$", req.content, flags=re.M)][:24] \
                or _TYPES.get(req.type, [])
        gen = {"content": req.content, "sections": _secs,
               "ai_provenance": {"served_by": "verbatim-ingest", "is_external": False,
                                 "note": "content supplied verbatim by the caller — not regenerated"}}
    else:
        gen = await _generate(req.type, req.title, req.brief, req.domain, req.vsb_id, req.sections)
    _dur = int((time.time() - _t0) * 1000)
    # Continual operational delivery within the LIVING QMS: every produced deliverable is gated by the
    # OWNED QMS (real, stateful), held to the §10 Solution-Quality Bar, recorded within the §8 organism.
    _owner = request_owner_id(user if isinstance(user, dict) else None, "default")
    qa = await assure_delivery(gen["content"], gen["sections"], label="deliverable", owner_id=_owner)
    # §13→§8: producing a living deliverable is real cognitive work — register it with the homeostatic
    # controller so it EXPENDS metabolic ATP and carries the organism posture (like the cognition paths).
    try:
        from agentic_core.ai.native.homeostasis import homeostasis
        homeo = homeostasis.assess(demand_nodes=len(gen["sections"]))
    except Exception:
        homeo = None
    now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    deliverable = {
        "id": f"deliv-{uuid.uuid4().hex[:8]}",
        # §14 (W320) — server-stamped ownership (never caller-claimed under auth)
        "owner_id": _owner,
        "type": req.type,
        "title": req.title or req.brief[:60],
        "brief": req.brief,
        "domain": req.domain,
        "vsb_id": req.vsb_id,
        "sections": gen["sections"],
        "content": gen["content"],
        "ai_provenance": gen["ai_provenance"],
        "quality_assurance": qa,
        "homeostasis": homeo,
        "versions": [{"brief": req.brief, "content": gen["content"],
                      "ai_provenance": gen["ai_provenance"], "quality_assurance": qa, "created_at": now}],
        "reusable": True, "rerunnable": True, "living": True,
        "created_at": now, "updated_at": now,
    }
    rows = _load()
    rows.append(deliverable)
    _save(rows)
    try:
        from agentic_core.organism.biobus import biobus
        biobus.fire_signal("motor", "deliverables.produce", f"{req.type}: {deliverable['title']}", 0.6)
    except Exception:
        pass
    try:
        from agentic_core.api.operational_excellence import record_outcome
        record_outcome("deliverable", f"deliverable:{req.type}",
                       served_by=gen["ai_provenance"]["served_by"],
                       is_external=gen["ai_provenance"]["is_external"],
                       duration_ms=_dur, success=bool(gen["content"]),
                       ref=deliverable["id"], vsb_id=req.vsb_id)
    except Exception:
        pass
    return deliverable


@router.get("")
async def list_deliverables(vsb_id: Optional[str] = None,
                            user: dict | None = Depends(get_current_user)):
    _u = user if isinstance(user, dict) else None
    rows = [d for d in _load() if user_can_access(_u, d.get("owner_id"))]   # §14 (W320)
    if vsb_id:
        rows = [d for d in rows if d.get("vsb_id") == vsb_id]
    summaries = [{"id": d["id"], "type": d["type"], "title": d["title"], "vsb_id": d.get("vsb_id"),
                  "versions": len(d.get("versions", [])), "served_by": d.get("ai_provenance", {}).get("served_by"),
                  "qms_gate_passed": d.get("quality_assurance", {}).get("quality", {}).get("qms_gate_passed"),
                  "updated_at": d.get("updated_at")} for d in rows]
    return {"deliverables": summaries[::-1], "total": len(summaries)}


@router.get("/{deliverable_id}")
async def get_deliverable(deliverable_id: str, user: dict | None = Depends(get_current_user)):
    for d in _load():
        if d["id"] == deliverable_id:
            _require_deliverable_access(d, user)
            return d
    raise HTTPException(status_code=404, detail=f"Deliverable {deliverable_id} not found.")


def _to_markdown(d: Dict[str, Any]) -> str:
    prov = d.get("ai_provenance", {})
    served = prov.get("served_by", "native")
    in_house = "in-house" if not prov.get("is_external") else f"via {served}"
    return (
        f"# {d.get('title', 'Deliverable')}\n\n"
        f"_{d.get('type', 'deliverable')} · produced on Workstation's own AI fabric ({in_house} · {served})_\n\n"
        f"> **Brief:** {d.get('brief', '')}\n\n"
        f"{d.get('content', '')}\n\n"
        "---\n"
        f"_Living deliverable `{d.get('id')}` — version {len(d.get('versions', []) or [1])}, "
        f"updated {d.get('updated_at', '')}. Re-generate or reconfigure it any time in Workstation IDBO._\n"
    )


# ── Honest in-house renderers (§4.9 / §13 — output in multiple selectable formats) ────────────────────
# Every format below is a REAL, deterministic in-house render of the deliverable's own content — no
# fabrication. Binary office/AV formats (pptx/docx/mp4/…) are NOT faked here; they remain a documented
# catalogue (see /output-formats) until the omnimedia generators produce them honestly.

def _md_inline(s: str) -> str:
    s = _htmlmod.escape(s)
    s = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)
    s = re.sub(r"`([^`]+)`", r"<code>\1</code>", s)
    s = re.sub(r"(?<!\*)\*(?!\s)([^*]+?)\*", r"<em>\1</em>", s)
    return s


def _md_to_html_body(md: str) -> str:
    """Minimal, faithful Markdown→HTML for the constructs deliverables use (headings, lists, blockquotes,
    rules, bold/italic/code, paragraphs). Deterministic structural render — not a guess."""
    out: List[str] = []
    in_ul = False
    for raw in md.split("\n"):
        line = raw.rstrip()
        if not line.strip():
            if in_ul:
                out.append("</ul>"); in_ul = False
            continue
        if line.strip() == "---":
            if in_ul:
                out.append("</ul>"); in_ul = False
            out.append("<hr/>"); continue
        h = re.match(r"^(#{1,6})\s+(.*)", line)
        if h:
            if in_ul:
                out.append("</ul>"); in_ul = False
            lvl = len(h.group(1)); out.append(f"<h{lvl}>{_md_inline(h.group(2))}</h{lvl}>"); continue
        li = re.match(r"^[-*]\s+(.*)", line)
        if li:
            if not in_ul:
                out.append("<ul>"); in_ul = True
            out.append(f"<li>{_md_inline(li.group(1))}</li>"); continue
        bq = re.match(r"^>\s?(.*)", line)
        if bq:
            if in_ul:
                out.append("</ul>"); in_ul = False
            out.append(f"<blockquote>{_md_inline(bq.group(1))}</blockquote>"); continue
        if in_ul:
            out.append("</ul>"); in_ul = False
        out.append(f"<p>{_md_inline(line)}</p>")
    if in_ul:
        out.append("</ul>")
    return "\n".join(out)


_DOC_CSS = (
    "body{font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;line-height:1.65;"
    "color:#1f2937;background:#f8fafc;margin:0;padding:2.5rem 1rem}main{max-width:820px;margin:0 auto;"
    "background:#fff;padding:3rem 3.5rem;border-radius:14px;box-shadow:0 1px 3px rgba(0,0,0,.08)}"
    "h1{font-size:2rem;letter-spacing:-.02em}h2{margin-top:2rem;border-bottom:1px solid #e5e7eb;padding-bottom:.3rem}"
    "blockquote{border-left:3px solid #14b8a6;margin:1rem 0;padding:.3rem 1rem;color:#475569;background:#f1f5f9}"
    "code{background:#f1f5f9;padding:.1rem .35rem;border-radius:4px;font-size:.9em}hr{border:none;border-top:1px solid #e5e7eb;margin:2rem 0}"
    ".sub{color:#64748b;font-size:.85rem;margin-top:-.5rem}.foot{color:#94a3b8;font-size:.75rem;margin-top:2.5rem}"
    "section.slide{page-break-after:always;border:1px solid #e5e7eb;border-radius:12px;padding:2rem;margin:1.5rem 0}"
)


def _doc_subtitle(d: Dict[str, Any]) -> str:
    prov = d.get("ai_provenance", {})
    served = prov.get("served_by", "native")
    mode = "in-house" if not prov.get("is_external") else f"via {served}"
    return f"{d.get('type', 'deliverable')} · produced on Workstation IDBO's own AI fabric ({mode} · {served})"


def _html_doc(d: Dict[str, Any]) -> str:
    title = d.get("title", "Deliverable")
    body = _md_to_html_body(d.get("content", ""))
    brief = d.get("brief", "")
    return (
        "<!doctype html><html lang=\"en\"><head><meta charset=\"utf-8\">"
        "<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">"
        f"<title>{_htmlmod.escape(title)}</title><style>{_DOC_CSS}</style></head><body><main>"
        f"<h1>{_htmlmod.escape(title)}</h1><p class=\"sub\">{_htmlmod.escape(_doc_subtitle(d))}</p>"
        + (f"<blockquote><strong>Brief:</strong> {_htmlmod.escape(brief)}</blockquote>" if brief else "")
        + body
        + f"<p class=\"foot\">Living deliverable {_htmlmod.escape(str(d.get('id','')))} — re-generate or reconfigure it any time in Workstation IDBO.</p>"
        "</main></body></html>"
    )


def _slides_doc(d: Dict[str, Any]) -> str:
    """A self-contained HTML slide deck — an honest 'Presentation' render: the deliverable split by its
    own ## section headings into one slide each (title slide first). Not a .pptx; a real, printable deck."""
    content = d.get("content", "")
    slides: List[tuple] = []
    cur_title, cur_lines = d.get("title", "Deliverable"), []
    for line in content.split("\n"):
        h = re.match(r"^##\s+(.*)", line)
        if h:
            slides.append((cur_title, "\n".join(cur_lines))); cur_title, cur_lines = h.group(1), []
        else:
            cur_lines.append(line)
    slides.append((cur_title, "\n".join(cur_lines)))
    body = [f"<section class=\"slide\"><h1>{_htmlmod.escape(d.get('title','Deliverable'))}</h1>"
            f"<p class=\"sub\">{_htmlmod.escape(_doc_subtitle(d))}</p></section>"]
    for stitle, sbody in slides:
        if not (stitle or sbody.strip()):
            continue
        body.append(f"<section class=\"slide\"><h2>{_htmlmod.escape(stitle)}</h2>{_md_to_html_body(sbody)}</section>")
    return (
        "<!doctype html><html lang=\"en\"><head><meta charset=\"utf-8\">"
        f"<title>{_htmlmod.escape(d.get('title','Deliverable'))} — slides</title><style>{_DOC_CSS}</style>"
        f"</head><body><main>{''.join(body)}</main></body></html>"
    )


def _video_html_doc(d: Dict[str, Any]) -> str:
    """A SELF-PLAYING video-style artifact — a real deterministic render (§4.9 'Video'): the
    deliverable's own sections become full-screen scenes that auto-advance with animated transitions
    and a progress bar (click to pause/resume, ←/→ to scrub, loops at the end). Self-contained HTML —
    no dependencies, honest content (nothing generated here beyond the deliverable itself). mp4/mp3
    remain in the not-yet catalogue until real media encoding exists."""
    content = d.get("content", "")
    scenes: List[tuple] = []
    cur_title, cur_lines = d.get("title", "Deliverable"), []
    for line in content.split("\n"):
        h = re.match(r"^##\s+(.*)", line)
        if h:
            scenes.append((cur_title, "\n".join(cur_lines))); cur_title, cur_lines = h.group(1), []
        else:
            cur_lines.append(line)
    scenes.append((cur_title, "\n".join(cur_lines)))
    body = [f"<section class=\"scene\"><h1>{_htmlmod.escape(d.get('title', 'Deliverable'))}</h1>"
            f"<p class=\"sub\">{_htmlmod.escape(_doc_subtitle(d))}</p></section>"]
    for stitle, sbody in scenes:
        if not (stitle or sbody.strip()):
            continue
        body.append(f"<section class=\"scene\"><h2>{_htmlmod.escape(stitle)}</h2>{_md_to_html_body(sbody)}</section>")
    css = (
        "html,body{margin:0;height:100%;background:#0b1220;color:#e6edf7;font-family:Segoe UI,system-ui,sans-serif}"
        ".scene{position:absolute;inset:0;display:flex;flex-direction:column;justify-content:center;"
        "padding:6vh 8vw;box-sizing:border-box;opacity:0;transform:translateY(24px);"
        "transition:opacity .8s ease,transform .8s ease;pointer-events:none;overflow:auto}"
        ".scene.active{opacity:1;transform:none;pointer-events:auto}"
        ".scene h1{font-size:3rem;margin:0 0 .5rem}.scene h2{font-size:2.2rem;margin:0 0 1rem;color:#8ab4ff}"
        ".scene .sub{color:#93a4c3}.scene p,.scene li{font-size:1.05rem;line-height:1.6;max-width:70ch}"
        "#bar{position:fixed;left:0;top:0;height:4px;background:#8ab4ff;width:0;transition:width .3s linear;z-index:9}"
        "#hint{position:fixed;right:1rem;bottom:.8rem;font-size:.7rem;color:#5b6b8c}"
    )
    js = (
        "const s=[...document.querySelectorAll('.scene')];let i=0,paused=false;const DUR=8000;"
        "const bar=document.getElementById('bar');"
        "function show(n){s.forEach((x,k)=>x.classList.toggle('active',k===n));"
        "bar.style.width=(100*(n+1)/s.length)+'%';}"
        "show(0);let t=setInterval(step,DUR);"
        "function step(){if(paused)return;i=(i+1)%s.length;show(i);}"
        "document.body.addEventListener('click',()=>{paused=!paused;});"
        "document.addEventListener('keydown',e=>{if(e.key==='ArrowRight'){i=(i+1)%s.length;show(i);}"
        "if(e.key==='ArrowLeft'){i=(i-1+s.length)%s.length;show(i);}});"
    )
    return (
        "<!doctype html><html lang=\"en\"><head><meta charset=\"utf-8\">"
        f"<title>{_htmlmod.escape(d.get('title', 'Deliverable'))} — video</title><style>{css}</style></head>"
        f"<body><div id=\"bar\"></div>{''.join(body)}"
        "<div id=\"hint\">auto-playing · click to pause · ←/→ to scrub</div>"
        f"<script>{js}</script></body></html>"
    )


def _strip_md(md: str) -> str:
    t = re.sub(r"^#{1,6}\s*", "", md, flags=re.M)
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"(?<!\*)\*(?!\s)([^*]+?)\*", r"\1", t)
    t = re.sub(r"(?<!_)_(?!\s)([^_]+?)_", r"\1", t)
    t = re.sub(r"`([^`]+)`", r"\1", t)
    t = re.sub(r"^>\s?", "", t, flags=re.M)
    t = re.sub(r"^[-*]\s+", "• ", t, flags=re.M)
    return t


# PDF is produced IN-HOUSE by the pure-python fpdf2 library when available (no external service, no
# system deps). Guarded so the app degrades gracefully if the lib is absent (PDF simply isn't offered).
try:
    import fpdf as _fpdf  # noqa: F401  (fpdf2 provides the `fpdf` module)
    _PDF_OK = True
except Exception:
    _PDF_OK = False

try:
    import docx as _docxmod  # noqa: F401  (python-docx → editable Word documents, pure-python)
    _DOCX_OK = True
except Exception:
    _DOCX_OK = False

try:
    import pptx as _pptxmod  # noqa: F401  (python-pptx → editable PowerPoint, pure-python)
    _PPTX_OK = True
except Exception:
    _PPTX_OK = False

try:
    import openpyxl as _openpyxlmod  # noqa: F401  (openpyxl → editable Excel, pure-python)
    _XLSX_OK = True
except Exception:
    _XLSX_OK = False


def _xml_safe(s: str) -> str:
    """Strip control characters that are illegal in OOXML/XML (so python-docx/openpyxl never raise on a
    deliverable whose content happens to contain a NULL byte or control char). Keeps tab/newline/CR."""
    return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", s or "")


def _demark(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    s = re.sub(r"(?<!\*)\*(?!\s)([^*]+?)\*", r"\1", s)
    return _xml_safe(s)


def _latin1(s: str) -> str:
    """fpdf2's core fonts are Latin-1; map the common Unicode punctuation our content uses to safe
    equivalents, then drop anything still unmappable. Honest: a real PDF, with minor glyph substitution."""
    for k, v in {"—": "-", "–": "-", "’": "'", "‘": "'", "“": '"',
                 "”": '"', "…": "...", "→": "->", "•": "-", " ": " "}.items():
        s = s.replace(k, v)
    return s.encode("latin-1", "replace").decode("latin-1")


def _pdf_bytes(d: Dict[str, Any]) -> bytes:
    """A real in-house PDF render of the deliverable (title · subtitle · brief · structured content)."""
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
    pdf = FPDF(format="A4")
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(20, 18, 20)
    pdf.add_page()

    def mc(h: float, text: str) -> None:
        # always return to the left margin on the next line (default advances X to the right margin,
        # which leaves zero width for the following cell — the cause of "no horizontal space").
        pdf.multi_cell(0, h, _latin1(text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.set_font("Helvetica", "B", 20); mc(9, d.get("title", "Deliverable"))
    pdf.set_text_color(110, 110, 110); pdf.set_font("Helvetica", "I", 9)
    mc(5, _doc_subtitle(d)); pdf.set_text_color(0, 0, 0); pdf.ln(2)
    if d.get("brief"):
        pdf.set_font("Helvetica", "I", 10); mc(5, "Brief: " + d["brief"]); pdf.ln(2)
    for raw in (d.get("content", "") or "").split("\n"):
        line = raw.rstrip()
        if not line.strip():
            pdf.ln(2); continue
        if line.strip() == "---":
            pdf.ln(1); continue
        h = re.match(r"^(#{1,6})\s+(.*)", line)
        if h:
            pdf.ln(1); pdf.set_font("Helvetica", "B", {1: 16, 2: 14, 3: 12}.get(len(h.group(1)), 11))
            mc(7, _demark(h.group(2))); continue
        li = re.match(r"^[-*]\s+(.*)", line)
        if li:
            pdf.set_font("Helvetica", "", 11); mc(6, "  -  " + _demark(li.group(1))); continue
        bq = re.match(r"^>\s?(.*)", line)
        if bq:
            pdf.set_text_color(80, 80, 80); pdf.set_font("Helvetica", "I", 11)
            mc(6, _demark(bq.group(1))); pdf.set_text_color(0, 0, 0); continue
        pdf.set_font("Helvetica", "", 11); mc(6, _demark(line))
    return bytes(pdf.output())


def _docx_bytes(d: Dict[str, Any]) -> bytes:
    """A real in-house editable Word (.docx) render of the deliverable (python-docx, pure-python).
    Unicode-native, so no glyph mapping needed; markdown structure → Word headings/bullets/quotes."""
    import io
    from docx import Document
    doc = Document()
    doc.add_heading(_xml_safe(d.get("title", "Deliverable")), level=0)
    doc.add_paragraph(_xml_safe(_doc_subtitle(d))).runs[0].italic = True
    if d.get("brief"):
        doc.add_paragraph(_xml_safe("Brief: " + d["brief"])).runs[0].italic = True
    for raw in (d.get("content", "") or "").split("\n"):
        line = raw.rstrip()
        if not line.strip() or line.strip() == "---":
            continue
        h = re.match(r"^(#{1,6})\s+(.*)", line)
        if h:
            doc.add_heading(_demark(h.group(2)), level=min(len(h.group(1)), 4)); continue
        li = re.match(r"^[-*]\s+(.*)", line)
        if li:
            doc.add_paragraph(_demark(li.group(1)), style="List Bullet"); continue
        bq = re.match(r"^>\s?(.*)", line)
        if bq:
            try:
                doc.add_paragraph(_demark(bq.group(1)), style="Quote")
            except Exception:
                doc.add_paragraph(_demark(bq.group(1))).runs[0].italic = True
            continue
        doc.add_paragraph(_demark(line))
    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


def _pptx_bytes(d: Dict[str, Any]) -> bytes:
    """A real in-house editable PowerPoint (.pptx) render — a title slide plus one content slide per
    deliverable ## section, bullets from the section body (python-pptx, pure-python, Unicode-native)."""
    import io
    from pptx import Presentation
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = _xml_safe(d.get("title", "Deliverable"))
    try:
        title_slide.placeholders[1].text = _xml_safe(_doc_subtitle(d))
    except Exception:
        pass
    # split the content into sections by its own ## headings (mirrors the HTML slide deck)
    sections: List[tuple] = []
    cur_title, cur_lines = None, []
    for line in (d.get("content", "") or "").split("\n"):
        h = re.match(r"^##\s+(.*)", line)
        if h:
            if cur_title is not None or cur_lines:
                sections.append((cur_title or d.get("title", ""), cur_lines))
            cur_title, cur_lines = h.group(1), []
        else:
            cur_lines.append(line)
    sections.append((cur_title or d.get("title", ""), cur_lines))
    for stitle, slines in sections:
        bullets = [_demark(re.sub(r"^\s*[#>*\-]+\s*", "", ln)).strip()
                   for ln in slines if ln.strip() and ln.strip() != "---"]
        if not (stitle or bullets):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = _xml_safe(str(stitle))[:140]
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets[:14]):   # cap bullets so a slide never overflows
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b[:200]
        except Exception:
            pass
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


def _xlsx_bytes(d: Dict[str, Any]) -> bytes:
    """A real in-house editable spreadsheet (.xlsx) render — a header plus a Section|Content table, one
    row per deliverable ## section (openpyxl, pure-python, Unicode-native)."""
    import io
    from openpyxl import Workbook
    from openpyxl.styles import Font
    wb = Workbook(); ws = wb.active; ws.title = "Deliverable"
    ws["A1"] = _xml_safe(d.get("title", "Deliverable")); ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = _xml_safe(_doc_subtitle(d))
    if d.get("brief"):
        ws["A3"] = _xml_safe("Brief: " + d["brief"])
    row = 5
    ws.cell(row, 1, "Section").font = Font(bold=True); ws.cell(row, 2, "Content").font = Font(bold=True)
    row += 1
    sections: List[tuple] = []
    cur_title, cur_lines = "Overview", []
    for line in (d.get("content", "") or "").split("\n"):
        h = re.match(r"^##\s+(.*)", line)
        if h:
            sections.append((cur_title, cur_lines)); cur_title, cur_lines = h.group(1), []
        else:
            cur_lines.append(line)
    sections.append((cur_title, cur_lines))
    for stitle, slines in sections:
        body = " ".join(_demark(re.sub(r"^\s*[#>*\-]+\s*", "", ln)).strip()
                        for ln in slines if ln.strip() and ln.strip() != "---")
        if not (stitle or body):
            continue
        ws.cell(row, 1, _xml_safe(str(stitle))[:240]); c = ws.cell(row, 2, body[:4000]); c.alignment = c.alignment.copy(wrap_text=True)
        row += 1
    ws.column_dimensions["A"].width = 30; ws.column_dimensions["B"].width = 96
    buf = io.BytesIO(); wb.save(buf); return buf.getvalue()


# Live, end-to-end in-house formats (real renders) — keep in sync with /output-formats.
# ── §13 (W372) — real vector + raster renders ────────────────────────────────────────────────
# mp4/mp3/png/svg were all catalogue-only. Two of the four can be produced HONESTLY in-house today
# and now are; the other two cannot, and stay catalogued rather than faked:
#   • svg — a text format, so a genuine self-contained vector render with NO dependency at all.
#   • png — a genuine raster render via Pillow when it is importable (optional, gated exactly like
#     pdf/docx/pptx/xlsx above).
#   • mp4/mp3 — need a real media encoder. There is no ffmpeg/av/imageio here, so they remain in the
#     not-yet catalogue. A "video" that is silently a slideshow of stills would be a fabrication.

_CARD_W, _CARD_H = 1200, 675          # 16:9 — shares and prints cleanly

# How many bullet lines actually FIT above the footer. Rendering the card and LOOKING at it showed
# the 8th line colliding with the provenance footer, and a repeated heading printing twice — both
# invisible when only checking that the bytes decode.
_CARD_MAX_LINES = 6


def _card_lines(d: Dict[str, Any]) -> List[str]:
    """The deliverable's OWN section headings — never invented filler. De-duplicated (a document
    that repeats a heading, or lists it in both `sections` and the body, must not print it twice)
    and capped to what fits above the footer."""
    seen = set()
    out: List[str] = []
    for line in (d.get("content", "") or "").split("\n"):
        m = re.match(r"^##+\s+(.*)", line)
        if m:
            head = m.group(1).strip()
            key = head.lower()
            if head and key not in seen:
                seen.add(key)
                out.append(head)
    if not out:
        for s in (d.get("sections") or []):
            if isinstance(s, str) and s.strip() and s.strip().lower() not in seen:
                seen.add(s.strip().lower())
                out.append(s.strip())
    return out[:_CARD_MAX_LINES]


def _svg_doc(d: Dict[str, Any]) -> str:
    """A real, self-contained SVG summary card of the deliverable (vector, no dependencies)."""
    import html as _h
    title = _h.escape((d.get("title") or "Deliverable")[:90])
    sub = _h.escape(_doc_subtitle(d)[:120])
    prov = d.get("ai_provenance") or {}
    served = _h.escape(str(prov.get("served_by") or "in-house"))
    posture = "external" if prov.get("is_external") else "in-house"
    rows = []
    y = 300
    for head in _card_lines(d):
        rows.append(
            f'<text x="90" y="{y}" font-family="Segoe UI, Helvetica, Arial, sans-serif" '
            f'font-size="30" fill="#cbd5e1">\u2022 {_h.escape(head[:70])}</text>')
        y += 46
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{_CARD_W}" height="{_CARD_H}" '
        f'viewBox="0 0 {_CARD_W} {_CARD_H}" role="img" aria-label="{title}">'
        f'<defs><linearGradient id="bg" x1="0" y1="0" x2="1" y2="1">'
        f'<stop offset="0%" stop-color="#0b1220"/><stop offset="100%" stop-color="#111c33"/>'
        f'</linearGradient></defs>'
        f'<rect width="{_CARD_W}" height="{_CARD_H}" fill="url(#bg)"/>'
        f'<rect x="0" y="0" width="{_CARD_W}" height="8" fill="#64ffda"/>'
        f'<text x="90" y="150" font-family="Segoe UI, Helvetica, Arial, sans-serif" font-size="56" '
        f'font-weight="700" fill="#ffffff">{title}</text>'
        f'<text x="90" y="205" font-family="Segoe UI, Helvetica, Arial, sans-serif" font-size="28" '
        f'fill="#94a3b8">{sub}</text>'
        + "".join(rows) +
        f'<text x="90" y="{_CARD_H - 60}" font-family="Segoe UI, Helvetica, Arial, sans-serif" '
        f'font-size="22" fill="#64ffda">Workstation IDBO \u00b7 {posture} \u00b7 {served}</text>'
        f'</svg>'
    )


try:
    from PIL import Image as _PILImage, ImageDraw as _PILDraw, ImageFont as _PILFont  # noqa: F401
    _PNG_OK = True
except Exception:
    _PNG_OK = False


def _png_bytes(d: Dict[str, Any]) -> bytes:
    """A real raster render of the same card via Pillow. Uses a real TrueType face when one can be
    found so the text is legible at size; falls back to Pillow's bitmap font rather than failing."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (_CARD_W, _CARD_H), (11, 18, 32))
    dr = ImageDraw.Draw(img)

    def _font(size: int):
        for cand in ("segoeui.ttf", "arial.ttf", "DejaVuSans.ttf",
                     "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"):
            try:
                return ImageFont.truetype(cand, size)
            except Exception:
                continue
        return ImageFont.load_default()

    dr.rectangle([0, 0, _CARD_W, 8], fill=(100, 255, 218))
    dr.text((90, 110), (d.get("title") or "Deliverable")[:70], font=_font(54), fill=(255, 255, 255))
    dr.text((90, 185), _doc_subtitle(d)[:110], font=_font(26), fill=(148, 163, 184))
    y = 280
    body = _font(28)
    for head in _card_lines(d):
        dr.text((90, y), f"\u2022 {head[:66]}", font=body, fill=(203, 213, 225))
        y += 44
    prov = d.get("ai_provenance") or {}
    posture = "external" if prov.get("is_external") else "in-house"
    dr.text((90, _CARD_H - 78), f"Workstation IDBO \u00b7 {posture} \u00b7 {prov.get('served_by') or 'in-house'}",
            font=_font(22), fill=(100, 255, 218))
    import io as _io
    buf = _io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


_LIVE_FORMATS = {
    "md":     ("text/markdown; charset=utf-8", "md"),
    "html":   ("text/html; charset=utf-8", "html"),
    "slides": ("text/html; charset=utf-8", "slides.html"),
    "txt":    ("text/plain; charset=utf-8", "txt"),
    "json":   ("application/json; charset=utf-8", "json"),
    "video-html": ("text/html; charset=utf-8", "video.html"),   # self-playing video-style render (W264)
    "svg":    ("image/svg+xml; charset=utf-8", "svg"),          # W372 — real vector card, no deps
}
if _PNG_OK:
    _LIVE_FORMATS["png"] = ("image/png", "png")                 # W372 — real raster card via Pillow
if _PDF_OK:
    _LIVE_FORMATS["pdf"] = ("application/pdf", "pdf")
if _DOCX_OK:
    _LIVE_FORMATS["docx"] = ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx")
if _PPTX_OK:
    _LIVE_FORMATS["pptx"] = ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx")
if _XLSX_OK:
    _LIVE_FORMATS["xlsx"] = ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx")


def _render_deliverable(d: Dict[str, Any], fmt: str) -> str:
    if fmt == "svg":
        return _svg_doc(d)
    if fmt == "md":
        return _to_markdown(d)
    if fmt == "txt":
        return _strip_md(_to_markdown(d))
    if fmt == "json":
        keep = ("id", "type", "title", "brief", "domain", "vsb_id", "sections", "content",
                "ai_provenance", "created_at", "updated_at")
        return json.dumps({k: d.get(k) for k in keep}, indent=2, ensure_ascii=False)
    if fmt == "html":
        return _html_doc(d)
    if fmt == "slides":
        return _slides_doc(d)
    if fmt == "video-html":
        return _video_html_doc(d)
    return _to_markdown(d)


@router.get("/{deliverable_id}/export")
async def export_deliverable(deliverable_id: str, format: str = "md",
                             user: dict | None = Depends(get_current_user)):
    """Export a living deliverable in a selectable in-house format (§4.9) so the user can take it out and
    use it. Live formats: md · html (styled document/website) · slides (HTML presentation) · txt · json —
    every one a REAL deterministic render of the deliverable's own content (no fabrication)."""
    fmt = (format or "md").lower()
    if fmt not in _LIVE_FORMATS:
        raise HTTPException(status_code=400,
                            detail=f"Unsupported format '{fmt}'. Live in-house formats: {sorted(_LIVE_FORMATS)}.")
    for d in _load():
        if d["id"] == deliverable_id:
            _require_deliverable_access(d, user)
            slug = "".join(c if c.isalnum() else "-" for c in d.get("title", "deliverable")).strip("-")[:60] or "deliverable"
            media_type, ext = _LIVE_FORMATS[fmt]
            if fmt == "pdf":
                content: Any = _pdf_bytes(d)
            elif fmt == "docx":
                content = _docx_bytes(d)
            elif fmt == "pptx":
                content = _pptx_bytes(d)
            elif fmt == "xlsx":
                content = _xlsx_bytes(d)
            elif fmt == "png":
                content = _png_bytes(d)
            else:
                content = _render_deliverable(d, fmt)
            return Response(
                content=content,
                media_type=media_type,
                headers={"Content-Disposition": f'attachment; filename="{slug}.{ext}"'},
            )
    raise HTTPException(status_code=404, detail=f"Deliverable {deliverable_id} not found.")


@router.post("/{deliverable_id}/regenerate")
async def regenerate(deliverable_id: str, req: RegenerateRequest,
                     user: dict | None = Depends(get_current_user)):
    """Re-run / reconfigure a living deliverable — appends a new version (history preserved)."""
    rows = _load()
    for d in rows:
        if d["id"] == deliverable_id:
            _require_deliverable_access(d, user)
            brief = req.brief or d["brief"]
            sections = req.sections if req.sections is not None else d.get("sections", [])
            if req.content.strip():
                # §3A (W308) — DEVELOP: a refined draft (e.g. /refine output) persists as the next
                # version VERBATIM — the develop loop lands in history instead of evaporating in the UI.
                gen = {"content": req.content, "sections": sections or [
                           m.group(1).strip() for m in
                           re.finditer(r"^#{1,2}\s+(.+)$", req.content, flags=re.M)][:24],
                       "ai_provenance": {"served_by": "verbatim-ingest", "is_external": False,
                                         "note": "refined content supplied verbatim — not regenerated"}}
            else:
                # §13 (W425) — hand the model the draft it is replacing. Skipped when the brief or
                # the section structure CHANGED: the user reconfigured the deliverable, so the prior
                # draft answers a different question and improving it would fight the instruction.
                _reconfigured = (brief != d["brief"]) or (sections != d.get("sections", []))
                _prior = "" if _reconfigured else (d.get("content") or "")
                gen = await _generate(d["type"], d["title"], brief, d.get("domain", "enterprise"),
                                      d.get("vsb_id"), sections, prior=_prior)
            qa = await assure_delivery(gen["content"], gen["sections"], label="deliverable",
                                       owner_id=d.get("owner_id"))
            now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

            # §13 (W425) — MEASURE whether this version is actually better, on the same real proxies
            # the QMS gate already computes. Previously a regeneration replaced `content`
            # unconditionally: a worse draft silently displaced a better one, and nothing anywhere
            # said the quality had gone down. The replacement still happens — the user asked for it —
            # but the verdict is recorded and returned so no one has to guess.
            _prev_q = ((d.get("quality_assurance") or {}).get("quality") or {})
            _new_q = (qa.get("quality") or {})
            _before = _prev_q.get("delivery_coverage")
            _after = _new_q.get("delivery_coverage")
            _delta = (round(_after - _before, 3)
                      if isinstance(_before, (int, float)) and isinstance(_after, (int, float))
                      else None)
            improvement = {
                # Taken from the generator's own report of the prompt it sent — never from whether
                # this endpoint intended to supply a prior draft.
                "compared_against_prior_draft": bool(gen.get("improved_from_prior")),
                "improvement_focus": gen.get("improvement_focus") or [],
                "coverage_before": _before,
                "coverage_after": _after,
                "coverage_delta": _delta,
                "sections_uncovered_before": _weakest_sections(d.get("content") or "", gen["sections"]),
                "sections_uncovered_after": _weakest_sections(gen["content"], gen["sections"]),
                "verdict": ("not_compared" if _delta is None else
                            "improved" if _delta > 0 else
                            "regressed" if _delta < 0 else "unchanged"),
                "note": ("This version scored LOWER than the one it replaced. The previous version is "
                         "preserved in history." if (_delta or 0) < 0 else None),
            }

            d["brief"] = brief
            d["sections"] = gen["sections"]
            d["content"] = gen["content"]
            d["ai_provenance"] = gen["ai_provenance"]
            d["quality_assurance"] = qa
            d["improvement"] = improvement
            d["updated_at"] = now
            d.setdefault("versions", []).append({"brief": brief, "content": gen["content"],
                                                 "ai_provenance": gen["ai_provenance"],
                                                 "quality_assurance": qa, "improvement": improvement,
                                                 "created_at": now})
            _save(rows)
            return d
    raise HTTPException(status_code=404, detail=f"Deliverable {deliverable_id} not found.")
